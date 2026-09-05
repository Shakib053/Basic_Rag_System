from __future__ import annotations

import csv
import hashlib
import mimetypes
import os
from datetime import datetime, timezone
from pathlib import Path
from uuid import NAMESPACE_URL, uuid5

from langchain_core.documents import Document


SUPPORTED_FILE_TYPES = {".txt", ".md", ".pdf", ".docx", ".pptx", ".html", ".htm", ".csv", ".xlsx"}
DEFAULT_MAX_UPLOAD_BYTES = 50 * 1024 * 1024
TABULAR_ROWS_PER_DOCUMENT = 100


class DocumentLoadError(ValueError):
    """Raised when an upload cannot be safely converted to searchable text."""


def max_upload_bytes() -> int:
    return int(os.getenv("MAX_UPLOAD_BYTES", str(DEFAULT_MAX_UPLOAD_BYTES)))


def stable_document_id(path: str | Path) -> str:
    canonical_path = str(Path(path).expanduser().resolve())
    return str(uuid5(NAMESPACE_URL, canonical_path))


def file_content_hash(path: str | Path) -> str:
    digest = hashlib.sha256()
    with Path(path).open("rb") as source:
        for block in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _base_metadata(path: Path, content_hash: str) -> dict:
    suffix = path.suffix.lower()
    return {
        "document_id": stable_document_id(path),
        "content_hash": content_hash,
        "file_name": path.name,
        "file_type": suffix.lstrip("."),
        "mime_type": mimetypes.guess_type(path.name)[0] or "application/octet-stream",
        "source": str(path.resolve()),
        "source_locator": str(path.resolve()),
        "ingested_at": datetime.now(timezone.utc).isoformat(),
    }


def validate_document_path(path: str | Path) -> Path:
    file_path = Path(path).expanduser().resolve()
    if not file_path.exists() or not file_path.is_file():
        raise DocumentLoadError(f"File does not exist: {file_path}")
    if file_path.suffix.lower() not in SUPPORTED_FILE_TYPES:
        supported = ", ".join(sorted(SUPPORTED_FILE_TYPES))
        raise DocumentLoadError(f"Unsupported file type '{file_path.suffix}'. Supported: {supported}.")
    size = file_path.stat().st_size
    if size == 0:
        raise DocumentLoadError("The file is empty.")
    if size > max_upload_bytes():
        limit_mb = max_upload_bytes() / (1024 * 1024)
        raise DocumentLoadError(f"File exceeds the configured {limit_mb:g} MiB upload limit.")
    return file_path


def _plain_text_documents(path: Path, metadata: dict) -> list[Document]:
    try:
        text = path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError as exc:
        raise DocumentLoadError("Text files must use UTF-8 encoding.") from exc
    return [Document(page_content=text, metadata=dict(metadata))] if text.strip() else []


def _pdf_documents(path: Path, metadata: dict) -> tuple[list[Document], list[str]]:
    import pymupdf

    warnings: list[str] = []
    try:
        pdf = pymupdf.open(path)
    except Exception as exc:
        raise DocumentLoadError(f"Could not open PDF: {exc}") from exc

    try:
        if pdf.needs_pass:
            raise DocumentLoadError("Encrypted PDFs are not supported.")
        documents = []
        blank_pages = 0
        for page_index, page in enumerate(pdf):
            text = page.get_text("text").strip()
            if not text:
                blank_pages += 1
                continue
            page_metadata = dict(metadata)
            page_metadata.update({"page": page_index, "source_locator": f"page {page_index + 1}"})
            documents.append(Document(page_content=text, metadata=page_metadata))
        if blank_pages:
            warnings.append(f"Skipped {blank_pages} page(s) with no extractable text.")
        return documents, warnings
    finally:
        pdf.close()


def _docx_documents(path: Path, metadata: dict) -> list[Document]:
    try:
        from docx import Document as WordDocument
    except ImportError as exc:
        raise DocumentLoadError("DOCX support requires the 'python-docx' package.") from exc

    try:
        word = WordDocument(path)
    except Exception as exc:
        raise DocumentLoadError(f"Could not open DOCX: {exc}") from exc

    blocks = [paragraph.text.strip() for paragraph in word.paragraphs if paragraph.text.strip()]
    for table in word.tables:
        rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        if rows:
            blocks.append("\n".join(rows))
    text = "\n\n".join(blocks)
    return [Document(page_content=text, metadata=dict(metadata))] if text.strip() else []


def _pptx_documents(path: Path, metadata: dict) -> list[Document]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentLoadError("PPTX support requires the 'python-pptx' package.") from exc

    try:
        presentation = Presentation(path)
    except Exception as exc:
        raise DocumentLoadError(f"Could not open PPTX: {exc}") from exc

    documents = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        blocks = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                blocks.append(text)
            if getattr(shape, "has_table", False):
                rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows]
                blocks.extend(rows)
        if blocks:
            slide_metadata = dict(metadata)
            slide_metadata.update({"slide": slide_index, "source_locator": f"slide {slide_index}"})
            documents.append(Document(page_content="\n\n".join(blocks), metadata=slide_metadata))
    return documents


def _html_documents(path: Path, metadata: dict) -> list[Document]:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise DocumentLoadError("HTML support requires the 'beautifulsoup4' package.") from exc

    try:
        markup = path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError as exc:
        raise DocumentLoadError("HTML files must use UTF-8 encoding.") from exc
    soup = BeautifulSoup(markup, "html.parser")
    for element in soup(["script", "style", "noscript"]):
        element.decompose()
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
    return [Document(page_content=text, metadata=dict(metadata))] if text else []


def _rows_to_documents(rows: list[list[object]], metadata: dict, sheet_name: str) -> list[Document]:
    if not rows:
        return []
    header = [str(value) if value is not None else "" for value in rows[0]]
    data_rows = rows[1:] if len(rows) > 1 else rows
    documents = []
    for offset in range(0, len(data_rows), TABULAR_ROWS_PER_DOCUMENT):
        batch = data_rows[offset : offset + TABULAR_ROWS_PER_DOCUMENT]
        lines = [" | ".join(header)] if len(rows) > 1 else []
        lines.extend(" | ".join(str(value) if value is not None else "" for value in row) for row in batch)
        text = "\n".join(line for line in lines if line.strip(" |"))
        if not text:
            continue
        row_start = offset + (2 if len(rows) > 1 else 1)
        row_end = row_start + len(batch) - 1
        row_metadata = dict(metadata)
        row_metadata.update(
            {
                "sheet_name": sheet_name,
                "row_start": row_start,
                "row_end": row_end,
                "source_locator": f"{sheet_name} rows {row_start}-{row_end}",
            }
        )
        documents.append(Document(page_content=text, metadata=row_metadata))
    return documents


def _csv_documents(path: Path, metadata: dict) -> list[Document]:
    try:
        with path.open("r", encoding="utf-8-sig", newline="") as source:
            rows = list(csv.reader(source))
    except UnicodeDecodeError as exc:
        raise DocumentLoadError("CSV files must use UTF-8 encoding.") from exc
    except csv.Error as exc:
        raise DocumentLoadError(f"Could not parse CSV: {exc}") from exc
    return _rows_to_documents(rows, metadata, "CSV")


def _xlsx_documents(path: Path, metadata: dict) -> list[Document]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise DocumentLoadError("XLSX support requires the 'openpyxl' package.") from exc
    try:
        workbook = load_workbook(path, read_only=True, data_only=True)
    except Exception as exc:
        raise DocumentLoadError(f"Could not open XLSX: {exc}") from exc
    try:
        documents = []
        for sheet in workbook.worksheets:
            documents.extend(_rows_to_documents(list(sheet.iter_rows(values_only=True)), metadata, sheet.title))
        return documents
    finally:
        workbook.close()


def load_document_file(path: str | Path) -> tuple[list[Document], list[str]]:
    """Load one validated file and return text documents plus non-fatal warnings."""
    file_path = validate_document_path(path)
    content_hash = file_content_hash(file_path)
    metadata = _base_metadata(file_path, content_hash)
    suffix = file_path.suffix.lower()
    warnings: list[str] = []

    if suffix in {".txt", ".md"}:
        documents = _plain_text_documents(file_path, metadata)
    elif suffix == ".pdf":
        documents, warnings = _pdf_documents(file_path, metadata)
    elif suffix == ".docx":
        documents = _docx_documents(file_path, metadata)
    elif suffix == ".pptx":
        documents = _pptx_documents(file_path, metadata)
    elif suffix in {".html", ".htm"}:
        documents = _html_documents(file_path, metadata)
    elif suffix == ".csv":
        documents = _csv_documents(file_path, metadata)
    else:
        documents = _xlsx_documents(file_path, metadata)

    if not documents:
        raise DocumentLoadError("The file contains no extractable text. OCR and vision are not enabled.")
    return documents, warnings


def load_documents_from_directory(data_dir: str | Path) -> tuple[list[Document], list[str]]:
    documents: list[Document] = []
    warnings: list[str] = []
    for path in sorted(Path(data_dir).iterdir(), key=lambda item: item.name.casefold()):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_FILE_TYPES:
            continue
        try:
            loaded, file_warnings = load_document_file(path)
        except DocumentLoadError as exc:
            warnings.append(f"{path.name}: {exc}")
            continue
        documents.extend(loaded)
        warnings.extend(f"{path.name}: {warning}" for warning in file_warnings)
    return documents, warnings
