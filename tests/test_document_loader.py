import csv
import tempfile
import unittest
from pathlib import Path

import pymupdf
from docx import Document as WordDocument
from openpyxl import Workbook
from pptx import Presentation

from ingestion.document_loader import DocumentLoadError, load_document_file, stable_document_id


class DocumentLoaderTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.root = Path(self.temp_dir.name)

    def tearDown(self):
        self.temp_dir.cleanup()

    def test_plain_text_markdown_html_and_csv(self):
        fixtures = {
            "note.txt": "Alpha fact",
            "readme.md": "# Beta\n\nBeta fact",
            "page.html": "<html><script>bad()</script><body><h1>Gamma</h1></body></html>",
        }
        for name, content in fixtures.items():
            path = self.root / name
            path.write_text(content, encoding="utf-8")
            documents, _ = load_document_file(path)
            self.assertTrue(documents[0].page_content.strip())
            self.assertEqual(documents[0].metadata["document_id"], stable_document_id(path))

        csv_path = self.root / "table.csv"
        with csv_path.open("w", encoding="utf-8", newline="") as target:
            csv.writer(target).writerows([["name", "value"], ["delta", 4]])
        documents, _ = load_document_file(csv_path)
        self.assertIn("delta | 4", documents[0].page_content)
        self.assertEqual(documents[0].metadata["source_locator"], "CSV rows 2-2")

    def test_docx_pptx_xlsx_and_pdf_preserve_locations(self):
        docx_path = self.root / "letter.docx"
        word = WordDocument()
        word.add_paragraph("Document fact")
        word.save(docx_path)
        documents, _ = load_document_file(docx_path)
        self.assertIn("Document fact", documents[0].page_content)

        pptx_path = self.root / "slides.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Slide fact"
        presentation.save(pptx_path)
        documents, _ = load_document_file(pptx_path)
        self.assertEqual(documents[0].metadata["slide"], 1)

        xlsx_path = self.root / "workbook.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Metrics"
        sheet.append(["name", "score"])
        sheet.append(["epsilon", 5])
        workbook.save(xlsx_path)
        documents, _ = load_document_file(xlsx_path)
        self.assertEqual(documents[0].metadata["source_locator"], "Metrics rows 2-2")

        pdf_path = self.root / "report.pdf"
        pdf = pymupdf.open()
        page = pdf.new_page()
        page.insert_text((72, 72), "PDF fact")
        pdf.save(pdf_path)
        pdf.close()
        documents, _ = load_document_file(pdf_path)
        self.assertEqual(documents[0].metadata["source_locator"], "page 1")

    def test_unsupported_empty_and_image_only_files_are_rejected(self):
        unsupported = self.root / "archive.zip"
        unsupported.write_bytes(b"not a zip")
        with self.assertRaisesRegex(DocumentLoadError, "Unsupported"):
            load_document_file(unsupported)

        empty = self.root / "empty.txt"
        empty.write_bytes(b"")
        with self.assertRaisesRegex(DocumentLoadError, "empty"):
            load_document_file(empty)

        scanned = self.root / "scanned.pdf"
        pdf = pymupdf.open()
        pdf.new_page()
        pdf.save(scanned)
        pdf.close()
        with self.assertRaisesRegex(DocumentLoadError, "no extractable text"):
            load_document_file(scanned)


if __name__ == "__main__":
    unittest.main()
